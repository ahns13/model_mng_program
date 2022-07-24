from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import copy
from model_sql_ppt_maker import *
from datetime import datetime


def modelPptMaker(**kargs):
    """
    :param kargs:
        model_keys: model key list
        model_names: model name list
        file_names: ppt file name list
        file_paths: ppt file location list of model

    :return:
    """
    shape_id = {
        "profile": 13, "career": 9, "contract": 10,
        "big image left": 3, "big image right": 2,
        "small image 1": 4, "small image 2": 5, "small image 3": 6, "small image 4": 7
    }
    prs = Presentation("./ppt/sample1.pptx")
    main_slide = prs.slides[0]
    model_size = len(kargs["model_names"])

    for i in range(model_size):
        if i > 0:
            add_slide = prs.slides.add_slide(main_slide.slide_layout)
            for shp in main_slide.shapes:
                if shp.shape_id in [9, 11]:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    add_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for i in range(model_size):
        model_profile_data = ppt_get_profile(kargs["model_keys"][i])
        model_career_data = ppt_get_career(kargs["model_keys"][i])
        model_contract_data = ppt_get_contract(kargs["model_keys"][i])
        print(model_profile_data)
        print(model_career_data)
        print(model_contract_data)
        for shp in prs.slides[i].shapes:
            if shp.shape_id == shape_id["profile"]:  # 모델 프로필 삽입
                tf = shp.text_frame
                pg = tf.paragraphs[0]  # name
                addRun(pg, " ".join(model_profile_data["name"]) + (" "+model_profile_data["name_info"] if model_profile_data["name_info"] else ""))
                if model_profile_data["birth_date"]:
                    pg = tf.add_paragraph()  # birth
                    pg.text = ""
                    addRun(pg, "생년월일 : " + model_profile_data["birth_date"])
                if model_profile_data["height"]:
                    pg = tf.add_paragraph()  # height
                    pg.text = ""
                    addRun(pg, "신장 : " + model_profile_data["height"])
                if model_profile_data["weight"]:
                    pg = tf.add_paragraph()  # weight
                    pg.text = ""
                    addRun(pg, "체중 : " + model_profile_data["weight"])
                if model_profile_data["sizes"]:
                    pg = tf.add_paragraph()  # sizes
                    pg.text = ""
                    addRun(pg, model_profile_data["sizes"])
                if model_profile_data["hobbynspec"]:
                    pg = tf.add_paragraph()  # hobbynspec
                    pg.text = ""
                    addRun(pg, model_profile_data["hobbynspec"])

            if shp.shape_id == shape_id["career"]:  # 모델 경력 삽입
                tf = shp.text_frame
                for idx, data in enumerate(model_career_data):
                    pg = tf.add_paragraph()  # career_type
                    pg.text = ""
                    addRun(pg, "- "+data[0])
                    pg = tf.add_paragraph()  # careers
                    pg.text = ""
                    addRun(pg, data[1])
                    if idx < len(model_career_data):
                        pg = tf.add_paragraph()  # blank line
                        pg.text = ""

            if shp.shape_id == shape_id["contract"]:  # 모델 경력 삽입
                tf = shp.text_frame
                for idx, data in enumerate(model_contract_data):
                    pg = tf.add_paragraph()  # career_type
                    pg.text = ""
                    pg.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    addRun(pg, data[2] + (" "+data[3] if data[3] else ""), "contract")
    cur_time = datetime.now().strftime("%Y%m%d")
    prs.save("./models_"+cur_time+".pptx" if model_size > 1 else "./"+kargs["model_names"][0]+".pptx")


def addRun(v_paragraph, v_run_text, v_type=None):
    run = v_paragraph.add_run()
    run.text = v_run_text
    run.font.name = "맑은 고딕"
    if v_type == "contract":
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(192, 0, 0)
    else:
        run.font.size = Pt(10)


if __name__ == "__main__":
    modelPptMaker(model_names=["test"], model_keys=[60])
