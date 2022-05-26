from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tkinter import filedialog
import re
from datetime import date
import sys, os
from program_info import *
import tkinter as tk
from model_sql import *

# 변수
re_car = r"[가-힣]+\s*:\s*"
re_s = re.compile("[0-9]+")
re_tel_text = "[0-9]{2,3}-[0-9]{3,4}-[0-9]{4}"  # 전화번호
re_name_text = r"([가-힣]+[a-z|A-Z|\s|\(|\)|:]*)+"
re_name_only = r"[가-힣|a-z|A-Z|\s]+"
re_name_desc = r"\([\w+(+|_|-|,)*\s*\w+]*\)$"
re_cont_name_n_tel = re.compile("(?:"+re_name_text+r"\s*)+"+re_tel_text)
re_cont_tel = re.compile(re_tel_text)
re_cont_date = re.compile(r"\(\d{6}\)")  # (등록일자(6))
re_chin_text = r"\([一-龥]\)"  # 본인 한자
re_insta = re.compile(r"([^\w+]@|^@)\s*[\w+(.|_|*)*]+")  # instram_id
re_amt_text = re.compile(r"[0-9]+[가-힣+\s?|\s+]+[0-9]{1}[~]*[0-9]{0,3}[.\d]*")  # 계약 금액
re_amt_ap_text = re.compile(r"\s*\(+[ap|AP]\w+[%|-]*\)+")

career_map = {
    # 입력단어: category
    "드라마": "드라마", "웹드라마": "웹드라마",
    "광고": "광고", "CF": "광고", "TVCF": "광고", "홍보물": "광고", "홍보": "광고", "AD": "광고", "홍보영상": "홍보영상",
    "뮤직비디오": "뮤직비디오", "M/V": "뮤직비디오", "뮤비": "뮤직비디오", "앨범": "앨범", "OST": "앨범", "음반": "앨범",
    "영화": "영화", "단편영화": "단편영화", "단편": "단편영화", "독립영화": "독립영화",
    "모델": "모델", "MODEL": "모델", "쇼호스트": "쇼호스트",
    "SHOW": "SHOW", "쇼": "SHOW", "패션쇼": "SHOW", "FASHION SHOW": "SHOW", "RUNWAY": "SHOW",
    "방송": "방송", "뮤지컬": "뮤지컬", "공연": "공연", "연극": "연극",
    "잡지": "잡지", "매거진": "잡지", "MEGAZINE": "잡지", "CA": "잡지", "카다로그": "카다로그",
    "LOOKBOOK": "LOOKBOOK", "룩북": "LOOKBOOK", "화보": "화보",
    "수상": "수상", "입상": "수상", "쇼핑몰": "쇼핑몰", "지면": "지면",
    "성우": "성우",
    "활동": "활동", "경력": "활동", "기타": "기타"
}

positions = ["사원", "주임", "대리", "과장", "차장", "본부장", "부장", "이사", "팀장", "실장", "상무", "전무", "부사장", "사장", "대표", "파트장", "사업부장"]
folder_file_data = {}


class GetOutLoop(Exception):
    pass

class DupCheck(Exception):
    pass

class Model:
    def __init__(self):
        self.file_excuted = False
        self.file_dup_bf_key = None
        self.model_profile = {
            "name": "",
            "real_name": "",
            "model_desc": "",
            "birthYear": "",
            "height": "",
            "weight": "",
            "sizeTop": "",
            "sizePants": "",
            "sizeShoe": "",
            "sizeOther": "",
            "hobbyNspec": [],
            "tel": "",
            "email": "",
            "insta_id": "",
            "data_date": ""
        }

        self.model_career = {
            "type": "",
            "드라마": {}, "웹드라마": {},
            "영화": {}, "단편영화": {}, "독립영화": {},
            "광고": {},
            "앨범": {}, "뮤직비디오": {},
            "모델": {}, "SHOW": {}, "쇼호스트": {},
            "방송": {},
            "뮤지컬": {}, "공연": {}, "연극": {},
            "잡지": {}, "카다로그": {}, "LOOKBOOK": {}, "화보": {}, "홍보영상": {},
            "수상": {}, "쇼핑몰": {}, "지면": {},
            "성우": {},
            "활동": {}, "기타": {}
        }

        self.model_contact = [
            # [이름 또는 소속사, 매니저/소속사 구분(M/C/""), 직급/직책, 주담당자 구분, 연락처, 이메일, 등록일자]
        ]

        self.model_contract_amt = {
            # contract_type : {'amount': [[month, amount, 등록일자]], 'ap': ""}
        }

        self.model_contract_info = [
            # [content, 등록일자]
        ]

        self.model_file_info = {
            "model_gubun": "",
            "model_dir_route": "",
            "file_name": "",
            "folder_path": ""
        }

    def career_divider(self, vGubun_name, vInput_list):
        for cont in vInput_list:
            if "웹" in cont:
                cont_ext = re.search(re.compile(r"\(\w*웹\w*\)"), cont).group()
                gubun_name_ext = re.sub(r"\(|\)", "", cont_ext).replace("웹", "")
                gubun_name_ext = gubun_name_ext if gubun_name_ext else vGubun_name
                cont_mod = re.search(re.compile(r"[\w+\s*\w*]+(?!=\(\w+\))"),  cont).group()
                if gubun_name_ext in self.model_career["웹드라마"]:
                    self.model_career["웹드라마"][gubun_name_ext].append(cont_mod)
                else:
                    self.model_career["웹드라마"][gubun_name_ext] = [cont_mod]
            else:
                if vGubun_name in self.model_career[self.model_career["type"]]:
                    self.model_career[self.model_career["type"]][vGubun_name].append(cont)
                else:
                    self.model_career[self.model_career["type"]][vGubun_name] = [cont]

    def file_exist_check(self, v_final_check):
        name_birth = self.model_profile["name"] + self.model_profile["birthYear"]
        if name_birth in folder_file_data["nameBirth"]:
            key_index = folder_file_data["nameBirth"].index(name_birth)
            if self.model_file_info["file_name"] in folder_file_data["fileName"]:
                return "e"  # exists
            elif v_final_check:
                if self.model_profile["tel"] != folder_file_data["tel"][key_index] and \
                   self.model_profile["insta_id"] != folder_file_data["instaId"][key_index]:
                    # 이름, 생년이 같은 동명이인인 경우 핸드폰번호와 인스타번호가 모두 다르면 새로 생성시킨다.
                    return "x"
                elif self.model_profile["tel"] == "" and self.model_profile["insta_id"] == "":
                    return "s"
                else:
                    return ["c", folder_file_data["key"][key_index]]
            else:
                return ["c", folder_file_data["key"][key_index]]  # check
        else:
            return "x"


def birthCompare(vBirth):
    cyear = str(date.today().year)
    if vBirth > cyear[2:] if len(vBirth) == 2 else cyear:
        return "19"+vBirth if len(vBirth) == 2 else vBirth
    else:
        return "20"+vBirth if len(vBirth) == 2 else vBirth


def data_ins_date(vRow_text):
    cont_date = re.search(re_cont_date, vRow_text)
    if cont_date is not None:
        return re.sub(r"\(|\)|c", "", cont_date.group())
    else:
        return ""


def position_check(vCont_text):
    position_name = ""
    for pos in positions:
        if pos in vCont_text:
            position_name = pos
            break
    return position_name


# shape left, top, height
model_profile_box = [133350, 549275, 792163]  # 신상정보
model_career_box = [133350, 1341438, 3602037]  # 경력사항
model_contact_box = [107950, 5013325, 1728788]  # 연락처및계약금액 : 박스 top,height가 다른 ppt 존재함


def shape_exec(cls_model, p_shape):
    shape_type = ""
    if model_profile_box[1] - 10000 < p_shape.top < model_profile_box[1] + 10**5 and model_profile_box[2] - 2*10**5 < p_shape.height < model_profile_box[2] + 2*10**5:  # profile_box
        shape_type = "A"
    elif model_career_box[1] - 10**5 < p_shape.top < model_career_box[1] + 10**5 and model_career_box[2] - 3*10**5 < p_shape.height:  # career_box
        shape_type = "B"
    elif model_contact_box[1] - 2*10**5 < p_shape.top < model_contact_box[1] + 6*10**5:  # contract_box
        shape_type = "C"

    gubun_name = "해당없음"  # 경력 상세 구분명
    category_in_chk = False
    career_ins_no = 0
    for p_idx, paraData in enumerate(p_shape.text_frame.paragraphs):
        row_text = paraData.text
        if shape_type == "A":
            data_ins_check = False
            if p_idx == 0:  # 이름
                name_text = row_text.replace(" ", "")
                if "(" in name_text and "本" not in name_text:
                    re_desc = re.search(re.compile(re_name_desc), row_text).group()
                    cls_model.model_profile["model_desc"] = re_desc
                    cls_model.model_profile["name"] = re.search(re.compile("[가-힣]+"), name_text).group()
                    data_ins_check = True
                elif re.search(re_s, name_text) is not None:
                    print(name_text)
                    name_text = re.search(re.compile("[가-힣]+"), name_text).group()

                if not data_ins_check:
                    if "本" in name_text:
                        cls_model.model_profile["name"] = name_text[:name_text.index("(本")]
                        cls_model.model_profile["real_name"] = name_text[name_text.index("(本")+2:len(name_text)-1]
                    else:
                        cls_model.model_profile["name"] = name_text
                    data_ins_check = True

            if "생년" in row_text or "년" in row_text:
                re_m = re.search(re_s, row_text)
                if re_m is not None:
                    cls_model.model_profile["birthYear"] = birthCompare(re_m.group())

                # data check
                check_val = cls_model.file_exist_check(False)
                if type(check_val) is list and check_val[0] == "c":
                    print("파일명은 다르나 이름이 존재합니다. 확인하세요.")
                    dup_ins_list = dup_ins_check(cls_model)
                    if cls_model.model_profile["name"] in dup_ins_list["name"] and \
                       cls_model.model_file_info["file_name"] in dup_ins_list["fileName"]:
                        cls_model.file_excuted = True
                        break
                    else:
                        pass
                elif check_val == "e":
                    print("현 파일은 작업된 파일입니다.")
                    cls_model.file_excuted = True
                    break
                data_ins_check = True
            elif "신장" in row_text or "cm" in row_text:
                re_m = re.search(re_s, row_text)
                if re_m is not None:
                    cls_model.model_profile["height"] = re_m.group()
                data_ins_check = True
            elif "체중" in row_text or "kg" in row_text:
                re_m = re.search(re_s, row_text)
                if re_m is not None:
                    cls_model.model_profile["weight"] = re_m.group()
                data_ins_check = True
            elif "특기" in row_text:
                re_m = re.sub("[특기|:| ]", "", row_text)
                if re_m != '':
                    cls_model.model_profile["hobbyNspec"].extend(re_m.split(","))
                data_ins_check = True
            elif "취미" in row_text:
                re_m = re.sub("[취미|:| ]", "", row_text)
                if re_m != '':
                    cls_model.model_profile["hobbyNspec"].extend(re_m.split(","))
                data_ins_check = True

            if "신발" in row_text or "슈즈" in row_text or "shoes" in row_text.lower():
                re_text = re.compile(r"[신발|슈즈|shoes|SHOES|Shoes][가-힣]*\s*:?\s*[0-9]{3}[~]*[0-9]{0,3}[mm]*")  # 신발사이즈: 000mm or 신발: 000
                re_m = re.search(re_text, row_text).group()
                m_num = re.search(re.compile("[0-9]{3}[~]*[0-9]{0,3}"), re_m).group()
                cls_model.model_profile["sizeShoe"] = m_num
                data_ins_check = True

            re_size_show = re.search(re.compile("[2|3][0-9]{2}(?!cm)"), row_text)
            if re_size_show is not None and int(re_size_show.group()) >= 220:  # 신발사이즈 : 신발 명시 없이 세자리 숫자(2xx)
                re_text = re.compile(r"\s*:?\s*[2|3][0-9]{2}[~]*[2|3]*[0-9]{0,2}[mm]*")  # 신발사이즈: 000mm or 신발: 000
                re_m = re.search(re_text, row_text).group()
                m_num = re.search(re.compile("[0-9]{3}[~]*[0-9]{0,3}"), re_m).group()
                cls_model.model_profile["sizeShoe"] = m_num
                data_ins_check = True

            re_top_pants = re.search(re.compile(r"상[.]*하[\w]*[_| ]*[0-9]{2}"), row_text)
            if re_top_pants is not None:
                cls_model.model_profile["size_other"] = re.search(re.compile("[0-9]{2}"), re_top_pants.group()).group()
                data_ins_check = True
            else:
                if "상의" in row_text:
                    re_text = re.compile(r"상의:?\s*[0-9]{2,3}[~]*[0-9]{0,3}")
                    re_m = re.search(re_text, row_text).group()
                    m_num = re.search(re.compile("[0-9]{2,3}[~]*[0-9]{0,3}"), re_m).group()
                    cls_model.model_profile["sizeTop"] = m_num
                    data_ins_check = True

                if "하의" in row_text:
                    re_text = re.compile(r"하의:?\s*[0-9]{2}[~]*[0-9]{0,2}")
                    re_m = re.search(re_text, row_text).group()
                    m_num = re.search(re.compile("[0-9]{2}[~]*[0-9]{0,2}"), re_m).group()
                    cls_model.model_profile["sizePants"] = m_num
                    data_ins_check = True

            re_size = re.search(re.compile("[0-9]{2}-[0-9]{2}-[0-9]{2}"), row_text.replace(" ", ""))
            if re_size is not None:
                cls_model.model_profile["sizeOther"] = re_size.group()
                data_ins_check = True

            if not data_ins_check:
                cls_model.model_profile["hobbyNspec"].extend(list(filter(None, row_text.replace(" ", "").split(","))))
        elif shape_type == "B":
            if "경력" in row_text.replace(" ", "") or not row_text:  # 경력사항 또는 빈 로우는 제외
                pass
            else:
                if re.search(re.compile("^[-|_]"), row_text) is not None:  # -로 시작하는 라인
                    category_name = re.sub("-|_| ", "", row_text).upper()
                    if len([1 for ctgy in career_map.keys() if ctgy == category_name]) > 0:
                        cls_model.model_career["type"] = career_map[category_name]
                    else:
                        cls_model.model_career["type"] = category_name

                        if len(p_shape.text_frame.paragraphs)-1 == p_idx or \
                           re.search(re.compile("^-"), p_shape.text_frame.paragraphs[p_idx+1].text) is not None:
                            if gubun_name in cls_model.model_career["기타"]:
                                cls_model.model_career["기타"][gubun_name].append(row_text.replace("-",  ""))
                            else:
                                cls_model.model_career["기타"][gubun_name] = [row_text.replace("-",  "")]
                        else:
                            cls_model.model_career[category_name] = {}

                    career_ins_no = 1
                    category_in_chk = True
                    gubun_name = "해당없음"
                elif category_in_chk:  # 경력 내용
                    # xxx : ~~~ -> 내용을 별도 구분지어 입력한 경우 [상세구분: DEFAIL_GUBUN]
                    cont_gubun_chk = re.search(re.compile(re_car), row_text)
                    if cont_gubun_chk is not None:
                        gubun_name = re.sub(re.compile(" |:"), "", cont_gubun_chk.group())
                        row_text = re.sub(re_car, "", row_text)
                        career_ins_no = 1
                    else:
                        gubun_name = gubun_name if gubun_name != "해당없음" else "해당없음"

                    career_text = re.sub(re.compile(", | ,| , "), ",", row_text).strip()

                    if re.search(re.compile(r"\(\w*웹\w*\)"), career_text) is not None:
                        input_list = list(filter(None, career_text.split(",")))
                        cls_model.career_divider(gubun_name, input_list)
                    else:
                        # (xx,x x,xx) 의 경우 별도 처리
                        career_list_check = re.search(re.compile(r"\([\w+,+\s*]+[^(]+\)"), career_text)
                        split_list = career_text.split(",")
                        input_list = [career_list_check.group()] if career_list_check is not None and "," in career_list_check.group() \
                            else list(filter(None, split_list))
                        if career_ins_no == 1:
                            cls_model.model_career[cls_model.model_career["type"]][gubun_name] = input_list
                            career_ins_no += 1
                        elif career_ins_no > 1:
                            cls_model.model_career[cls_model.model_career["type"]][gubun_name].extend(input_list)
                            career_ins_no += 1
                else:  # 카테고리 구분 없는 경력
                    if gubun_name in cls_model.model_career["기타"]:
                        cls_model.model_career["기타"][gubun_name].append(row_text)
                    else:
                        cls_model.model_career["기타"][gubun_name] = [row_text]
        elif shape_type == "C":
            if "Enter" in row_text or "Seoul" in row_text:
                pass
            else:
                data_ins_chk = False
                if re.search(re_cont_name_n_tel, row_text) is not None:  # 소속사 또는 매니저 + 연락처
                    # 소속사 (등록일자) 형식 처리 x -> 비고
                    cont_name_n_tel = re.search(re_cont_name_n_tel, row_text).group()
                    if "★" in row_text:
                        manage_gubun = "MAIN"
                    elif "광고)" in row_text:
                        manage_gubun = "광고"
                        cont_name_n_tel = re.sub(r"광고\)\s*", "", cont_name_n_tel)
                    else:
                        manage_gubun = ""

                    cont_name = re.search(re.compile(re_name_only), cont_name_n_tel).group().strip()
                    cont_tel = re.search(re_cont_tel, cont_name_n_tel).group()
                    position = position_check(cont_name)
                    manager_gubun = "N"
                    if position:  # 직책, 직급이 이름에 존재하면 이를 추출
                        cont_name = re.sub(" |"+position, "", cont_name)
                        manager_gubun = "M"

                    cls_model.model_contact.append([cont_name, manager_gubun, position, manage_gubun, cont_tel, "", data_ins_date(row_text)])
                    data_ins_chk = True
                elif re.search(re.compile(r"^\s*" + re_tel_text), row_text) is not None:  # 본인 연락처
                    cont_tel_text = re.search(re.compile(r"^\s*" + re_tel_text), row_text).group()
                    cont_tel = re.search(re_cont_tel, cont_tel_text).group()
                    cls_model.model_profile["tel"] = cont_tel
                    cls_model.model_profile["data_date"] = data_ins_date(row_text)
                    data_ins_chk = True
                elif re.search(re.compile(r"^\s*"+re_chin_text+r"\s*"+re_tel_text), row_text) is not None:  # (本) + 연락처
                    cont_tel_text = re.search(re.compile(r"^\s*"+re_chin_text+r"\s*"+re_tel_text), row_text).group()
                    cont_tel = re.search(re_cont_tel, cont_tel_text).group()
                    cls_model.model_profile["tel"] = cont_tel
                    cls_model.model_profile["data_date"] = data_ins_date(row_text)
                    data_ins_chk = True

                # 이메일 추출??

                if re.search(re_insta, row_text) is not None:  # instragram
                    insta_id = re.search(re_insta, row_text)
                    cls_model.model_profile["insta_id"] = insta_id[0].replace(" ", "") if insta_id else ""
                    data_ins_chk = True

                re_tel_chk = re.search(re_tel_text, row_text)
                if re_tel_chk is None and "@" not in row_text and ("개월" in row_text or len(re.findall(re.compile("[3|6|1]"), row_text))):
                    amt_type = re.search(re.compile(r"^[\(*가-힣+\)*\s*]+"), row_text)
                    if amt_type is not None:
                        model_amt_type = re.sub(r" ", "", amt_type.group())
                    else:
                        model_amt_type = "해당없음"

                    re_amt = re.findall(re_amt_text, row_text)
                    if len(re_amt):
                        cls_model.model_contract_amt[model_amt_type] = {'amount': []}
                        for amt in re_amt:
                            month_amt = re.findall(re.compile(r"[0-9][~]*[0-9]{0,3}[.\d]*"), amt)
                            if len(month_amt):
                                # if month_amt[0] not in ["3", "6", "1"]:
                                #     print("계약금액 추출 중 3,6,1 이외의 숫자 추출됨")
                                #     sys.exit()
                                # else:
                                cls_model.model_contract_amt[model_amt_type]['amount'].append([
                                    ("12" if month_amt[0] == "1" else month_amt[0]), float(month_amt[1])*10000 if "." in month_amt[1] else month_amt[1], data_ins_date(row_text)
                                ])
                                data_ins_chk = True

                        re_amt_ap = re.search(re_amt_ap_text, row_text)
                        cls_model.model_contract_amt[model_amt_type]['ap'] = re.sub(r" |\(|\)", "", re_amt_ap.group()) if re_amt_ap is not None else ""

                if not data_ins_chk and len(row_text):
                    content_ext = re.sub(r"s*\(+\d{6}\)+\s*", "", row_text)
                    cls_model.model_contract_info.append([content_ext, data_ins_date(row_text)])


def ppt_file_exec(vFilePath):
    if "~" in os.path.basename(vFilePath):
        return
    prs = Presentation(vFilePath)
    print('---- ', vFilePath)
    model = Model()
    model.model_file_info["model_gubun"] = os.path.basename(os.path.dirname(vFilePath))
    model.model_file_info["model_dir_route"] = os.path.dirname(vFilePath).replace(os.path.abspath(access_root), "")
    model.model_file_info["file_name"] = os.path.basename(vFilePath)
    model.model_file_info["folder_path"] = vFilePath

    global folder_file_data
    if not folder_file_data:
        folder_file_data = get_file_list(model.model_file_info["model_dir_route"])

    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for g_idx, g_shape in enumerate(shape.shapes):
                        if g_shape.left < 200000 and g_shape.height > 300000:  # Profile shape와 빈 shape 제외
                            if hasattr(g_shape, "text") and len(g_shape.text):
                                shape_exec(model, g_shape)
                                if model.file_excuted == "dup_check":
                                    raise DupCheck
                                elif model.file_excuted:
                                    raise GetOutLoop
                elif shape.has_text_frame and len(shape.text):
                    if shape.left < 200000 and shape.height > 300000:
                        shape_exec(model, shape)
                        if model.file_excuted == "dup_check":
                            raise DupCheck
                        elif model.file_excuted:
                            raise GetOutLoop

        check_val = model.file_exist_check(True)
        if type(check_val) is list and check_val[0] == "c":
            model.file_excuted = "dup_check"
            model.file_dup_bf_key = check_val[1]
            raise DupCheck
        elif check_val == "s":
            db_close()
            raise Exception("중복 이름(동명이인) 확인 필요!")
    except GetOutLoop:
        model.model_career["type"] = ""
        model.file_excuted = False
        return
    except DupCheck:
        model.model_career["type"] = ""
        model.file_excuted = False
        model_dup_ins(model)
        print("** 중복 파일로 추가됨")
        conn.commit()
        return

    model.model_career["type"] = ""
    print(model.model_profile)
    print(model.model_career)
    print(model.model_contact)
    print(model.model_contract_amt)
    print(model.model_contract_info)
    print(model.model_file_info)

    key_value = get_key()
    commit_status = []
    commit_status.append(model_profile_ins(model, key_value))  # insert model_profile
    commit_status.append(hoobynspec_ins(model, key_value))  # insert hobbynspec
    commit_status.append(career_ins(model, key_value))  # insert career
    commit_status.append(contact_ins(model, key_value))  # insert contact
    commit_status.append(contract_amount_ins(model, key_value))  # insert contract
    commit_status.append(contract_info_ins(model, key_value))  # insert contract_info
    commit_status.append(file_info_ins(model, key_value))  # insert file_info
    if commit_status:
        conn.commit()
        folder_file_data["key"].append(key_value)
        folder_file_data["nameBirth"].append(model.model_profile["name"]+model.model_profile["birthYear"])
        folder_file_data["fileName"].append(model.model_file_info["file_name"])
    else:
        db_cancel()
        inbox.destroy()
        sys.exit()


# folder_file_compare_db 실행 후 결과 list를 copy
mod_list = ['구종근92_2.pptx', '김우래83.pptx']


def ppt_data_exec(vType):
    if vType == "file":
        ppt_file = None
        while ppt_file is None:
            ppt_file = filedialog.askopenfilename(initialdir=access_root, filetypes=[("ppt files", "*.pptx;*.ppt")])
        if not len(mod_list) or os.path.basename(ppt_file) in mod_list:
            ppt_file_exec(os.path.abspath(ppt_file))
        else:
            print(os.path.basename(ppt_file)+" 파일이 mod_list 안에 존재하지 않습니다.")
        db_close()
        inbox.destroy()
    else:
        ppt_folder = None
        while ppt_folder is None:
            ppt_folder = filedialog.askdirectory(initialdir=access_root)

        max_file_data = get_max_data(os.path.abspath(ppt_folder).replace(os.path.abspath(access_root), ""))
        search_check = False
        for file in sorted(os.scandir(ppt_folder), key=lambda x: x.name):
            if file.is_file() and (file.name.endswith(".ppt") or file.name.endswith(".pptx")):
                if not len(mod_list):
                    if file.name == max_file_data:
                        search_check = True
                        continue
                    if max_file_data == "" or search_check:
                        ppt_file_exec(os.path.abspath(file.path))
                elif file.name in mod_list:
                    ppt_file_exec(os.path.abspath(file.path))
        db_close()
        inbox.destroy()


inbox = tk.Tk()
inbox.geometry("275x220")
inbox.title("ppt to db")
btn = tk.Button(inbox, text="단일파일", width=20, command=lambda: ppt_data_exec("file"))
btn.place(x=50, y=30)
btn2 = tk.Button(inbox, text="폴더", width=20, command=lambda: ppt_data_exec("folder"))
btn2.place(x=50, y=75)
inbox.lift()
inbox.focus_force()
# inbox.attributes("-topmost", True)
inbox.mainloop()
